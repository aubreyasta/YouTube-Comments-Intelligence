"""assets.py - text extraction from uploaded files and article fetching.

Both public functions are defensive: they never raise on bad input and
always return the documented type, so a run can continue past bad assets.
"""

import ipaddress
import logging
import os
import socket
import time
from datetime import datetime, timezone
from urllib.parse import urljoin, urlsplit, urlunsplit

logger = logging.getLogger(__name__)

_TEXT_EXTENSIONS = {".pdf", ".docx", ".pptx"}
_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp"}


class BlockedUrl(Exception):
    """An article URL that can reach a non-public address. Caller returns 422."""


_MAX_REDIRECTS = 5
_FETCH_BUDGET_SECONDS = 15.0
_DEFAULT_PORTS = {"http": 80, "https": 443}


def _is_public_ip(raw: str) -> bool:
    """True only for a globally routable unicast address.

    Rejects loopback, private, link-local, multicast, reserved, unspecified,
    IPv4-mapped IPv6 wrapping a non-public IPv4, and IPv6 unique-local.
    """
    try:
        ip = ipaddress.ip_address(raw)
    except ValueError:
        return False

    mapped = getattr(ip, "ipv4_mapped", None)
    if mapped is not None:
        ip = mapped

    if ip.version == 6 and ip in ipaddress.ip_network("fc00::/7"):
        return False

    return not (
        ip.is_loopback
        or ip.is_private
        or ip.is_link_local
        or ip.is_multicast
        or ip.is_reserved
        or ip.is_unspecified
    )


def _resolve_public_ips(host: str, port: int) -> list:
    """Resolve host and return its addresses, or raise BlockedUrl.

    A mixed public/non-public answer is rejected outright rather than
    filtered: an attacker who controls DNS can order the answers, so
    trusting the public half of a mixed reply trusts the attacker.
    """
    try:
        infos = socket.getaddrinfo(host, port, proto=socket.IPPROTO_TCP)
    except OSError:
        raise BlockedUrl("hostname does not resolve")

    addresses = [info[4][0] for info in infos]
    if not addresses:
        raise BlockedUrl("hostname resolved to no addresses")

    for address in addresses:
        if not _is_public_ip(address):
            raise BlockedUrl("hostname resolves to a non-public address")

    return addresses


def _validate_url(url: str) -> tuple:
    """Validate one absolute URL and resolve it. Raises BlockedUrl.

    Returns (scheme, host, port, addresses).
    """
    try:
        parts = urlsplit(url)
    except ValueError:
        raise BlockedUrl("malformed URL")

    scheme = (parts.scheme or "").lower()
    if scheme not in _DEFAULT_PORTS:
        raise BlockedUrl("scheme is not http or https")

    if parts.username or parts.password:
        raise BlockedUrl("URL carries embedded credentials")

    if parts.fragment:
        raise BlockedUrl("URL carries a fragment")

    try:
        host = parts.hostname
        port = parts.port
    except ValueError:
        raise BlockedUrl("malformed host or port")

    if not host:
        raise BlockedUrl("URL has no hostname")

    if port is not None and port != _DEFAULT_PORTS[scheme]:
        raise BlockedUrl("URL uses a non-default port")

    addresses = _resolve_public_ips(host, _DEFAULT_PORTS[scheme])
    return scheme, host, _DEFAULT_PORTS[scheme], addresses


def extract_upload(file_path: str) -> str:
    """Extract text from an uploaded file. Returns "" on images or unknown types.

    Supported:
      .pdf   - pypdf; returns "" and logs a warning on encrypted/empty PDFs.
      .docx  - python-docx; joins paragraph text.
      .pptx  - python-pptx; joins shape text across slides (no speaker notes).
      .png/.jpg/.jpeg/.webp - returns "" (OCR deferred to 12.1).
      other  - returns "".

    Never raises.
    """
    ext = os.path.splitext(file_path)[1].lower()

    if ext in _IMAGE_EXTENSIONS:
        return ""

    if ext not in _TEXT_EXTENSIONS:
        return ""

    try:
        if ext == ".pdf":
            return _extract_pdf(file_path)
        if ext == ".docx":
            return _extract_docx(file_path)
        if ext == ".pptx":
            return _extract_pptx(file_path)
    except Exception as exc:  # noqa: BLE001
        logger.warning("extract_upload failed for %s: %s", file_path, exc)
        return ""

    return ""


def _extract_pdf(file_path: str) -> str:
    import pypdf  # local import keeps module importable without pypdf installed

    try:
        reader = pypdf.PdfReader(file_path)
    except pypdf.errors.FileNotDecryptedError:
        logger.warning("PDF is encrypted, cannot extract text: %s", file_path)
        return ""

    pages = []
    for page in reader.pages:
        try:
            text = page.extract_text() or ""
            pages.append(text)
        except Exception as exc:  # noqa: BLE001
            logger.warning("PDF page extraction failed in %s: %s", file_path, exc)

    result = "\n".join(pages).strip()
    if not result:
        logger.warning("PDF yielded no text (possibly scanned): %s", file_path)
    return result


def _extract_docx(file_path: str) -> str:
    import docx as python_docx

    doc = python_docx.Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs if p.text)


def _extract_pptx(file_path: str) -> str:
    from pptx import Presentation

    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line:
                        texts.append(line)
    return "\n".join(texts)


def fetch_article(url: str) -> dict:
    """Fetch and extract readable text from a public URL.

    Returns a dict with keys:
      title       (str)  - page <title> or "" on failure
      text        (str)  - extracted body text, capped at 20 000 chars; "" on failure
      retrieved_at (str) - ISO-8601 UTC timestamp of the attempt

    Raises BlockedUrl when the URL - or any redirect hop - can reach a
    non-public address. Every other failure (timeout, network error,
    unparseable body) returns empty text, so a run continues past a bad
    asset (see Risks table: "Article fetch hangs").

    The request is sent to the address this function resolved, with the
    original hostname preserved for TLS (`sni_hostname`) and the `Host`
    header. Resolving here and connecting to the same answer closes the
    DNS-rebinding gap a validate-then-request-by-name design leaves open.
    """
    import httpx
    from bs4 import BeautifulSoup

    retrieved_at = datetime.now(timezone.utc).isoformat()
    empty = {"title": "", "text": "", "retrieved_at": retrieved_at}

    deadline = time.monotonic() + _FETCH_BUDGET_SECONDS
    current = url
    response = None

    try:
        with httpx.Client(follow_redirects=False) as client:
            for _ in range(_MAX_REDIRECTS + 1):
                scheme, host, port, addresses = _validate_url(current)

                remaining = deadline - time.monotonic()
                if remaining <= 0:
                    logger.warning("fetch_article timed out for %s", url)
                    return empty

                parts = urlsplit(current)
                literal = addresses[0]
                netloc = f"[{literal}]" if ":" in literal else literal
                pinned = urlunsplit((scheme, netloc, parts.path, parts.query, ""))

                response = client.get(
                    pinned,
                    timeout=remaining,
                    headers={"Host": host},
                    extensions={"sni_hostname": host},
                )

                if response.is_redirect:
                    location = response.headers.get("location")
                    if not location:
                        break
                    current = urljoin(current, location)
                    continue

                response.raise_for_status()
                break
            else:
                raise BlockedUrl("too many redirects")

        if response is None or response.is_redirect:
            return empty

    except BlockedUrl:
        raise
    except Exception as exc:  # noqa: BLE001
        logger.warning("fetch_article failed for %s: %s", url, exc)
        return empty

    try:
        soup = BeautifulSoup(response.text, "html.parser")

        # Remove noise before extracting text
        for tag in soup(["script", "style"]):
            tag.decompose()

        title = (soup.title.string or "").strip() if soup.title else ""
        text = soup.get_text(separator="\n", strip=True)[:20_000]

        return {"title": title, "text": text, "retrieved_at": retrieved_at}
    except Exception as exc:  # noqa: BLE001
        logger.warning("fetch_article parse failed for %s: %s", url, exc)
        return empty
